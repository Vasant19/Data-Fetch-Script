import streamlit as st
import mysql.connector
from mysql.connector import Error
import lida
from lida import Manager, TextGenerationConfig, llm
import base64
from PIL import Image
from io import BytesIO
from dotenv import load_dotenv
import os
import openai
import pandas as pd
import numpy as np
import random
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# Load environment variables from .env file (e.g., API keys)
load_dotenv()

# Retrieve the OpenAI API key from environment variables
api_key = os.getenv("OPENAI_API_KEY")

# Check if the API key is found; if not, raise an error
if api_key:
    openai.api_key = api_key
else:
    raise ValueError("OpenAI API key not found in environment variables!")

# Database connection string (Update this if not using MySQL database, such as snowflake etc)
CONNECTION_STRING = {
    "host": "127.0.0.1",  # Localhost
    "user": "vasant",  # Database user
    "password": "12345"  # Database password
}

# Helper function to convert base64-encoded string into an image object
def base64_to_image(base64_string):
    byte_data = base64.b64decode(base64_string)  # Decode the base64 string to byte data
    return Image.open(BytesIO(byte_data))  # Convert byte data to an image object using PIL

# Helper function to set a random seed for reproducibility
def set_random_seed(seed):
    np.random.seed(seed)  # Set numpy random seed
    random.seed(seed)  # Set Python random seed

# Function to generate a PowerPoint file with an image
def generate_ppt(image):
    prs = Presentation()  # Create a new PowerPoint presentation
    
    # Add a blank slide to the presentation (Layout 5 is typically a blank slide)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Convert the image to a stream (in-memory)
    img_stream = BytesIO()
    image.save(img_stream, format='PNG')  # Save the image as PNG into the in-memory stream
    img_stream.seek(0)  # Reset the stream to the start
    
    # Add the image to the slide at a given position (0.5 inches from top-left corner) with a specified size
    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(0.5), width=Inches(9))
    
    # Save the PowerPoint presentation to an in-memory buffer
    ppt_stream = BytesIO()
    prs.save(ppt_stream)  # Save the presentation to the stream
    ppt_stream.seek(0)  # Reset the stream to the start
    
    return ppt_stream  # Return the in-memory PowerPoint file

# Initialize session state variables to manage application state across interactions
if 'connection_established' not in st.session_state:
    st.session_state.connection_established = False  # Track if the database connection is established
if 'selected_database' not in st.session_state:
    st.session_state.selected_database = None  # Track selected database
if 'selected_table' not in st.session_state:
    st.session_state.selected_table = None  # Track selected table
if 'chart_image' not in st.session_state:
    st.session_state.chart_image = None  # Track generated chart image

# Streamlit UI setup (Title of the application)
st.title("Database Table Viewer")

# Step 1: Connect to the MySQL database when the user clicks the button
if st.button("Step 1: Connect to Database"):
    try:
        st.write("Attempting to connect to the database...")  # Inform the user about the connection attempt
        connection = mysql.connector.connect(**CONNECTION_STRING)  # Try to establish a connection to MySQL
        
        if connection.is_connected():  # Check if the connection was successful
            st.success("Connection to the database was successful!")  # Notify user of success
            db_info = connection.get_server_info()  # Get information about the MySQL server
            st.write(f"Connected to MySQL Server version {db_info}")  # Display server version
            st.session_state.connection_established = True  # Store connection state in session
            st.session_state.connection = connection  # Store the connection object for future use
        else:
            st.error("Failed to connect to the database")  # Notify user of connection failure

    except Error as e:
        st.error(f"Connection failed: {str(e)}")  # Catch and display any connection errors

# Step 2: Fetch and allow the user to select a database after a successful connection
if st.session_state.connection_established:
    try:
        cursor = st.session_state.connection.cursor()  # Create a cursor to execute SQL queries
        cursor.execute("SHOW DATABASES")  # Fetch all available databases
        databases = cursor.fetchall()  # Store the list of databases
        databases_list = [db[0] for db in databases]  # Extract database names into a list

        # Provide a dropdown for the user to select a database
        selected_db = st.selectbox(
            "Select a database",
            databases_list,
            key='database_selector',
            index=databases_list.index(st.session_state.selected_database)
            if st.session_state.selected_database in databases_list else 0
        )

        # If the selected database is different from the current one, update the session state
        if selected_db != st.session_state.selected_database:
            st.session_state.selected_database = selected_db  # Update the selected database
            st.session_state.selected_table = None  # Reset selected table
            st.rerun()  # Re-run the script to reflect the change

        cursor.close()  # Close the cursor

    except Error as e:
        st.error(f"Failed to fetch databases: {str(e)}")  # Catch and display any errors while fetching databases

# Step 3: Fetch and allow the user to select a table from the selected database
if st.session_state.selected_database:
    try:
        db_connection = mysql.connector.connect(
            **CONNECTION_STRING,
            database=st.session_state.selected_database  # Connect to the selected database
        )
        cursor = db_connection.cursor()  # Create a cursor to execute SQL queries
        cursor.execute("SHOW TABLES")  # Fetch all available tables in the selected database
        tables = cursor.fetchall()  # Store the list of tables
        tables_list = [table[0] for table in tables]  # Extract table names into a list

        # If tables exist, allow the user to select one
        if tables_list:
            selected_table = st.selectbox(
                "Select a table",
                tables_list,
                key='table_selector',
                index=tables_list.index(st.session_state.selected_table)
                if st.session_state.selected_table in tables_list else 0
            )

            # If the selected table is different from the current one, update the session state
            if selected_table != st.session_state.selected_table:
                st.session_state.selected_table = selected_table  # Update the selected table
                st.rerun()  # Re-run the script to reflect the change
        else:
            st.warning("No tables found in the selected database")  # If no tables are found, warn the user

        cursor.close()  # Close the cursor
        db_connection.close()  # Close the database connection

    except Error as e:
        st.error(f"Failed to fetch tables: {str(e)}")  # Catch and display any errors while fetching tables

# Step 4: View the data from the selected table
if st.session_state.selected_table:
    if st.button("View Table Data"):
        try:
            query = f"SELECT * FROM {st.session_state.selected_table}"  # Construct the SQL query to fetch all rows from the table
            db_connection = mysql.connector.connect(
                **CONNECTION_STRING,
                database=st.session_state.selected_database  # Connect to the selected database
            )
            df = pd.read_sql(query, db_connection)  # Fetch the data into a pandas DataFrame
            st.session_state.df = df  # Store the DataFrame in session state

            # Display the first few rows of the DataFrame
            if not df.empty:
                st.write(df.head())  # Display the first few rows of data
            else:
                st.info("No data found in the selected table")  # If the table is empty, notify the user

            db_connection.close()  # Close the database connection

        except Error as e:
            st.error(f"Failed to fetch data: {str(e)}")  # Catch and display any errors while fetching data

# Display current database and table selections in the sidebar
if st.session_state.selected_database:
    st.sidebar.write(f"Current Database: {st.session_state.selected_database}")
if st.session_state.selected_table:
    st.sidebar.write(f"Current Table: {st.session_state.selected_table}")

# Final Step: Generate visualization based on user input
if 'selected_table' in st.session_state and 'df' in st.session_state:
    try:
        # Initialize LIDA Manager for text generation using OpenAI
        lida = Manager(text_gen=llm("openai"))
        textgen_config = TextGenerationConfig(
            n=1, temperature=0.5, model="gpt-4o-mini", use_cache=True
        )

        # User input for generating visualizations based on the selected table
        user_query = st.text_input(
            "Enter your query for visualization:",
            placeholder="E.g., Give me a simple chart relevant to the selected table. Make sure the labelled text is visible "
        )

        if st.button("Generate Visualization"):
            set_random_seed(42)  # Set random seed for reproducibility
            if not user_query.strip():  # Check if user query is empty
                st.warning("Please enter a valid query!")  # Warn if the input is empty
            else:
                try:
                    df = st.session_state.df  # Get the DataFrame from session state
                    summary = lida.summarize(data=df, summary_method="default", textgen_config=textgen_config)  # Summarize data
                    st.session_state.summary = summary  # Store summary in session state

                    # Visualize the summarized data based on the user's query
                    library = "seaborn"  # Specify the library for visualizations
                    charts = lida.visualize(
                        summary=summary,
                        goal=user_query,
                        textgen_config=textgen_config,
                        library=library
                    )

                    # Display the generated chart
                    if charts and len(charts) > 0:
                        st.write(f"Number of charts generated: {len(charts)}")
                        image_base64 = charts[0].raster  # Get the base64-encoded image of the first chart
                        img = base64_to_image(image_base64)  # Convert to image object
                        st.session_state.chart_image = img  # Store the image in session state
                        st.image(img)  # Display the image
                    else:
                        st.error("No charts generated.")  # Notify if no charts are generated
                except Exception as e:
                    st.error(f"Failed to generate visualization: {str(e)}")  # Catch and display any errors

    except Exception as e:
        st.error(f"Error initializing LIDA Manager or configuration: {str(e)}")  # Catch errors related to LIDA setup

# Export the generated chart to a PowerPoint file
if st.session_state.chart_image:
    if st.button("Export to PowerPoint"):
        ppt_stream = generate_ppt(st.session_state.chart_image)  # Generate PowerPoint file with chart
        
        # Provide the PowerPoint file for download
        st.download_button(
            label="Download PowerPoint",
            data=ppt_stream,
            file_name="generated_chart.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
